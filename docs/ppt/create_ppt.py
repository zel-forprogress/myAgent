"""
myAgent 技术文档 PPT 生成脚本
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ============ 颜色主题 ============
BG_DARK = RGBColor(0x0F, 0x19, 0x29)        # 深蓝黑背景
BG_CARD = RGBColor(0x1A, 0x2B, 0x45)        # 卡片背景
ACCENT = RGBColor(0x00, 0xD4, 0xFF)         # 青色强调
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)        # 紫色辅助
ACCENT3 = RGBColor(0x10, 0xB9, 0x81)        # 绿色辅助
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)           # 灰色文字
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)     # 浅灰
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
PINK = RGBColor(0xEC, 0x48, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

blank_layout = prs.slide_layouts[6]


def set_solid_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_solid_fill(shp, fill)
    if not line:
        shp.line.fill.background()
    return shp


def add_rounded_rect(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_solid_fill(shp, fill)
    shp.line.fill.background()
    shp.adjustments[0] = 0.06
    return shp


def add_text(slide, x, y, w, h, text, font_size=14, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return tb


def add_circle(slide, cx, cy, r, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    set_solid_fill(shp, fill)
    shp.line.fill.background()
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT, width=2):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    # arrow head
    line_elem = line._element
    ln = line_elem.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(line_elem, qn('a:ln'))
    tailEnd = ln.find(qn('a:tailEnd'))
    if tailEnd is None:
        tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle')
    return line


def add_code_block(slide, x, y, w, h, code_lines, font_size=10, color=LIGHT_GRAY):
    """Add a code block with dark background"""
    add_rect(slide, x, y, w, h, RGBColor(0x0D, 0x11, 0x17))
    add_rect(slide, x, y, w, Pt(24), RGBColor(0x16, 0x1B, 0x22))
    # line numbers area
    add_rect(slide, x, Pt(24), Pt(30), h - Pt(24), RGBColor(0x0D, 0x11, 0x17))
    tb = slide.shapes.add_textbox(x + Pt(36), y + Pt(6), w - Pt(42), h - Pt(12))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Consolas"
        run.font.name = "Courier New"
    return tb


def add_page_header(slide, title, page_num=None, total=None):
    """Add consistent header bar"""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.6), BG_CARD)
    add_rect(slide, 0, Inches(0.6), SLIDE_W, Pt(2), ACCENT)
    add_text(slide, Inches(0.4), Inches(0.05), Inches(10), Inches(0.5),
             title, font_size=18, bold=True)
    if page_num and total:
        add_text(slide, Inches(11.5), Inches(0.05), Inches(1.5), Inches(0.5),
                 f"{page_num} / {total}", font_size=12, color=GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_bg(slide):
    """Set slide background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


# ============================================================
# Slide 1: Cover
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# Decorative elements
add_rect(slide, 0, 0, SLIDE_W, Pt(4), ACCENT)
add_circle(slide, Inches(11), Inches(1.5), Inches(1.2), ACCENT2)
add_circle(slide, Inches(12.5), Inches(2.5), Inches(0.6), ACCENT3)
add_circle(slide, Inches(2), Inches(6.5), Inches(0.8), RGBColor(0x1E, 0x3A, 0x5F))

# Title
add_text(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(1.2),
         "myAgent", font_size=52, bold=True, color=ACCENT)
add_text(slide, Inches(0.8), Inches(2.7), Inches(10), Inches(0.8),
         "本地 RAG Agent 系统技术文档", font_size=28, color=WHITE)

# Subtitle
add_rect(slide, Inches(0.8), Inches(3.7), Inches(3), Pt(3), ACCENT2)
add_text(slide, Inches(0.9), Inches(3.55), Inches(4), Inches(0.5),
         "FastAPI · LangGraph · Milvus · Next.js", font_size=14, bold=True)

# Tags
tags = ["RAG", "Agent", "多知识库", "会话记忆", "可观测性"]
tag_x = Inches(0.8)
for tag in tags:
    w = Inches(len(tag) * 0.18 + 0.5)
    add_rounded_rect(slide, tag_x, Inches(4.3), w, Inches(0.35), BG_CARD)
    add_text(slide, tag_x, Inches(4.3), w, Inches(0.35), tag, font_size=11,
             color=ACCENT, alignment=PP_ALIGN.CENTER)
    tag_x += w + Inches(0.15)

# Bottom info
add_text(slide, Inches(0.8), Inches(6.3), Inches(6), Inches(0.4),
         "技术栈：Python · FastAPI · LangGraph · Milvus · PostgreSQL · Redis · Celery · Next.js",
         font_size=11, color=GRAY)
add_text(slide, Inches(0.8), Inches(6.7), Inches(6), Inches(0.4),
         "模型：Qwen / DashScope OpenAI-compatible API",
         font_size=11, color=GRAY)


# ============================================================
# Slide 2: 项目背景
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_page_header(slide, "01  项目背景", 2, 12)

# Left: motivation
add_rounded_rect(slide, Inches(0.4), Inches(1.0), Inches(5.8), Inches(6.0), BG_CARD)
add_rect(slide, Inches(0.4), Inches(1.0), Inches(0.08), Inches(6.0), ACCENT)
add_text(slide, Inches(0.7), Inches(1.2), Inches(5.2), Inches(0.4),
         "为什么做 myAgent？", font_size=18, bold=True, color=ACCENT)

motivations = [
    ("大模型落地痛点", [
        "企业文档分散，LLM 无法直接利用私有知识",
        "通用对话缺少领域上下文，回答不精准",
        "长对话 token 成本高，需要会话记忆优化",
    ]),
    ("RAG 挑战", [
        "检索质量不稳定，需要改写和重排策略",
        "文档入库耗时，需要异步 Pipeline",
        "Agent 执行链路黑盒，需要可观测性",
    ]),
]

y_off = Inches(1.8)
for title, items in motivations:
    add_text(slide, Inches(0.7), y_off, Inches(5.2), Inches(0.35),
             title, font_size=14, bold=True, color=WHITE)
    y_off += Inches(0.4)
    for item in items:
        add_text(slide, Inches(0.9), y_off, Inches(5.0), Inches(0.3),
                 f"▸  {item}", font_size=11, color=LIGHT_GRAY)
        y_off += Inches(0.35)
    y_off += Inches(0.2)

# Right: goals
add_rounded_rect(side := slide, Inches(6.8), Inches(1.0), Inches(6.1), Inches(6.0), BG_CARD)
add_rect(slide, Inches(6.8), Inches(1.0), Inches(0.08), Inches(6.0), ACCENT2)
add_text(slide, Inches(7.1), Inches(1.2), Inches(5.5), Inches(0.4),
         "项目目标", font_size=18, bold=True, color=ACCENT2)

goals = [
    ("本地优先", "所有依赖 Docker Compose 一键启动，数据完全本地可控"),
    ("Agentic RAG", "基于 LangGraph 的多节点工作流，支持意图路由和检索质量检查"),
    ("多知识库", "支持创建多个独立知识库，检索时灵活选择"),
    ("会话记忆", "长对话自动摘要压缩，降低 token 消耗"),
    ("异步入库", "Redis + Celery 异步处理，支持重试和取消"),
    ("全链路观测", "Langfuse 记录每个 Agent 节点、检索结果和模型调用"),
]

y_off = Inches(1.8)
for i, (title, desc) in enumerate(goals):
    # number circle
    cx = Inches(7.3)
    cy = y_off + Inches(0.15)
    add_circle(slide, cx, cy, Inches(0.15), ACCENT3)
    add_text(slide, cx - Inches(0.15), cy - Inches(0.15), Inches(0.3), Inches(0.3),
             str(i + 1), font_size=10, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.6), y_off, Inches(2), Inches(0.35),
             title, font_size=13, bold=True, color=WHITE)
    add_text(slide, Inches(7.6), y_off + Inches(0.32), Inches(5.0), Inches(0.3),
             desc, font_size=10, color=GRAY)
    y_off += Inches(0.85)


# ============================================================
# Slide 3: 技术栈总览
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_page_header(slide, "02  技术栈总览", 3, 12)

# Title
add_text(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "全栈技术选型", font_size=24, bold=True, color=WHITE)

# Tech categories
categories = [
    {
        "title": "后端服务",
        "color": ACCENT,
        "items": ["Python 3.11", "FastAPI", "Uvicorn", "Pydantic v2"],
        "x": Inches(0.4), "y": Inches(1.7)
    },
    {
        "title": "Agent 编排",
        "color": ACCENT2,
        "items": ["LangGraph", "LangChain", "OpenAI-compatible API"],
        "x": Inches(3.5), "y": Inches(1.7)
    },
    {
        "title": "数据存储",
        "color": ACCENT3,
        "items": ["Milvus 2.5", "PostgreSQL 16", "MinIO", "Redis 7"],
        "x": Inches(6.6), "y": Inches(1.7)
    },
    {
        "title": "异步任务",
        "color": ORANGE,
        "items": ["Celery", "Redis Broker", "Pipeline 节点日志"],
        "x": Inches(9.7), "y": Inches(1.7)
    },
]

for cat in categories:
    x, y = cat["x"], cat["y"]
    w, h = Inches(2.9), Inches(2.2)
    add_rounded_rect(slide, x, y, w, h, BG_CARD)
    add_rect(slide, x, y, w, Pt(3), cat["color"])
    add_text(slide, x + Inches(0.2), y + Inches(0.1), w - Inches(0.4), Inches(0.35),
             cat["title"], font_size=14, bold=True, color=cat["color"])
    for i, item in enumerate(cat["items"]):
        add_text(slide, x + Inches(0.3), y + Inches(0.55) + Inches(0.35) * i,
                 w - Inches(0.5), Inches(0.3),
                 f"• {item}", font_size=12, color=LIGHT_GRAY)

# Second row
categories2 = [
    {
        "title": "文档解析",
        "color": PINK,
        "items": ["Docling", "Azure Document Intelligence", "pypdf", "python-docx"],
        "x": Inches(0.4), "y": Inches(4.2)
    },
    {
        "title": "可观测性",
        "color": ACCENT,
        "items": ["Langfuse", "链路追踪", "Token 用量统计"],
        "x": Inches(3.5), "y": Inches(4.2)
    },
    {
        "title": "前端",
        "color": ACCENT2,
        "items": ["Next.js 16", "React 19", "TypeScript", "Radix UI"],
        "x": Inches(6.6), "y": Inches(4.2)
    },
    {
        "title": "部署",
        "color": ACCENT3,
        "items": ["Docker Compose", "健康检查", "数据卷持久化"],
        "x": Inches(9.7), "y": Inches(4.2)
    },
]

for cat in categories2:
    x, y = cat["x"], cat["y"]
    w, h = Inches(2.9), Inches(2.2)
    add_rounded_rect(slide, x, y, w, h, BG_CARD)
    add_rect(slide, x, y, w, Pt(3), cat["color"])
    add_text(slide, x + Inches(0.2), y + Inches(0.1), w - Inches(0.4), Inches(0.35),
             cat["title"], font_size=14, bold=True, color=cat["color"])
    for i, item in enumerate(cat["items"]):
        add_text(slide, x + Inches(0.3), y + Inches(0.55) + Inches(0.35) * i,
                 w - Inches(0.5), Inches(0.3),
                 f"• {item}", font_size=12, color=LIGHT_GRAY)


# ============================================================
# Slide 4: 系统架构图
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_page_header(slide, "03  系统架构", 4, 12)

add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
         "端到端数据流与组件交互", font_size=20, bold=True, color=WHITE)

# --- Architecture diagram ---
# Layer 1: Users
add_rounded_rect(slide, Inches(0.3), Inches(1.4), Inches(1.8), Inches(0.6), BG_CARD)
add_text(slide, Inches(0.3), Inches(1.4), Inches(1.8), Inches(0.6),
         "👤 用户", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Layer 2: Frontend
add_rounded_rect(slide, Inches(2.5), Inches(1.3), Inches(2.2), Inches(0.8), ACCENT2)
add_text(slide, Inches(2.5), Inches(1.3), Inches(2.2), Inches(0.35),
         "Next.js 前端", font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(2.5), Inches(1.65), Inches(2.2), Inches(0.35),
         "聊天 · 登录 · 管理后台", font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow: user -> frontend
add_arrow(slide, Inches(2.1), Inches(1.7), Inches(2.5), Inches(1.7), ACCENT, 1.5)

# Layer 3: Backend (FastAPI)
add_rounded_rect(slide, Inches(5.2), Inches(1.2), Inches(2.8), Inches(1.0), ACCENT)
add_text(slide, Inches(5.2), Inches(1.2), Inches(2.8), Inches(0.35),
         "FastAPI 后端", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.2), Inches(1.55), Inches(2.8), Inches(0.3),
         "Auth · Chat · KB · Ingest", font_size=9, color=RGBColor(0xB0, 0xE0, 0xFF),
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.2), Inches(1.85), Inches(2.8), Inches(0.3),
         "Sessions · Retrieval · Admin", font_size=9, color=RGBColor(0xB0, 0xE0, 0xFF),
         alignment=PP_ALIGN.CENTER)

# Arrow: frontend -> backend
add_arrow(slide, Inches(4.7), Inches(1.7), Inches(5.2), Inches(1.7), ACCENT, 1.5)

# Layer 4: Agent (LangGraph)
add_rounded_rect(slide, Inches(8.5), Inches(1.2), Inches(2.2), Inches(1.0), ACCENT3)
add_text(slide, Inches(8.5), Inches(1.2), Inches(2.2), Inches(0.35),
         "LangGraph Agent", font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.5), Inches(1.55), Inches(2.2), Inches(0.3),
         "意图路由 · 检索", font_size=9, color=RGBColor(0xA0, 0xE0, 0xC0),
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.5), Inches(1.85), Inches(2.2), Inches(0.3),
         "改写 · 质量检查", font_size=9, color=RGBColor(0xA0, 0xE0, 0xC0),
         alignment=PP_ALIGN.CENTER)

# Arrow: backend -> agent
add_arrow(slide, Inches(8.0), Inches(1.7), Inches(8.5), Inches(1.7), ACCENT3, 1.5)

# Layer 5: LLM
add_rounded_rect(slide, Inches(11.1), Inches(1.3), Inches(1.9), Inches(0.8), ORANGE)
add_text(slide, Inches(11.1), Inches(1.3), Inches(1.9), Inches(0.35),
         "Qwen LLM", font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(11.1), Inches(1.65), Inches(1.9), Inches(0.35),
         "DashScope API", font_size=9, color=RGBColor(0xFF, 0xE0, 0xB0),
         alignment=PP_ALIGN.CENTER)

# Arrow: agent -> LLM
add_arrow(slide, Inches(10.7), Inches(1.7), Inches(11.1), Inches(1.7), ORANGE, 1.5)

# --- Bottom layer: Storage & Infrastructure ---
# Box for storage
add_rect(slide, Inches(0.3), Inches(2.8), Inches(12.7), Inches(0.02), RGBColor(0x2A, 0x3B, 0x55))

add_text(slide, Inches(0.5), Inches(2.9), Inches(3), Inches(0.35),
         "存储与基础设施", font_size=14, bold=True, color=ACCENT)

# Storage boxes
storages = [
    ("PostgreSQL", "用户 · 会话 · 消息\n知识库元数据 · 任务", ACCENT, Inches(0.3), Inches(3.4)),
    ("Milvus", "向量存储\n语义检索", ACCENT2, Inches(3.4), Inches(3.4)),
    ("MinIO", "原始文件存储\nS3 兼容", ACCENT3, Inches(6.0), Inches(3.4)),
    ("Redis", "Celery Broker\n任务队列", ORANGE, Inches(8.6), Inches(3.4)),
    ("Langfuse", "链路追踪\nToken 用量", PINK, Inches(11.2), Inches(3.4)),
]

for name, desc, color, x, y in storages:
    w = Inches(2.3)
    h = Inches(1.0)
    add_rounded_rect(slide, x, y, w, h, BG_CARD)
    add_rect(slide, x, y, Inches(0.06), h, color)
    add_text(slide, x + Inches(0.2), y + Inches(0.08), w - Inches(0.3), Inches(0.3),
             name, font_size=12, bold=True, color=color)
    add_text(slide, x + Inches(0.2), y + Inches(0.38), w - Inches(0.3), Inches(0.55),
             desc, font_size=9, color=GRAY)

# --- Bottom: Async Pipeline ---
add_rect(slide, Inches(0.3), Inches(4.7), Inches(12.7), Inches(0.02), RGBColor(0x2A, 0x3B, 0x55))
add_text(slide, Inches(0.5), Inches(4.8), Inches(3), Inches(0.35),
         "异步入库 Pipeline", font_size=14, bold=True, color=ORANGE)

pipeline_steps = [
    ("1. 上传文档", ACCENT),
    ("2. 创建 Task", ACCENT2),
    ("3. Celery 排队", ACCENT3),
    ("4. inspect_document", ORANGE),
    ("5. chunk_embed_index", PINK),
    ("6. update_record", ACCENT),
]

px = Inches(0.3)
for step_name, step_color in pipeline_steps:
    w = Inches(1.9)
    h = Inches(0.6)
    add_rounded_rect(slide, px, Inches(5.3), w, h, BG_CARD)
    add_rect(slide, px, Inches(5.3), Inches(0.06), h, step_color)
    add_text(slide, px, Inches(5.3), w, h, step_name, font_size=10,
             color=WHITE, alignment=PP_ALIGN.CENTER)
    # arrow to next
    if px < Inches(11):
        add_arrow(slide, px + w, Inches(5.6), px + w + Inches(0.12), Inches(5.6), GRAY, 1)
    px += w + Inches(0.12)

# --- Key features at bottom ---
add_rect(slide, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.02), RGBColor(0x2A, 0x3B, 0x55))
add_text(slide, Inches(0.5), Inches(6.4), Inches(12), Inches(0.35),
         "核心能力：多知识库检索  ·  会话记忆(摘要压缩)  ·  上下文补全  ·  检索改写  ·  Rerank  ·  流式输出",
         font_size=11, color=GRAY)


# ============================================================
# Slide 5: 核心流程 - 聊天问答
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_page_header(slide, "04  核心流程 — 聊天问答", 5, 12)

add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
         "LangGraph Agent 工作流", font_size=20, bold=True, color=WHITE)

# Flow diagram - nodes
nodes = [
    ("START", ACCENT, Inches(0.3), Inches(1.5)),
    ("complete_question\n_with_history", ACCENT2, Inches(2.0), Inches(1.5)),
    ("analyze_question\n(意图路由)", ACCENT2, Inches(4.0), Inches(1.5)),
    ("plan_agent_task", ACCENT2, Inches(6.0), Inches(1.5)),
]

# Draw top flow
for i, (name, color, x, y) in enumerate(nodes):
    w = Inches(1.6)
    h = Inches(0.8)
    add_rounded_rect(slide, x, y, w, h, BG_CARD)
    add_rect(slide, x, y, w, Pt(2), color)
    add_text(slide, x, y + Inches(0.05), w, h, name, font_size=9,
             color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < len(nodes) - 1:
        add_arrow(slide, x + w, y + Inches(0.4), x + w + Inches(0.4), y + Inches(0.4), color, 1.5)

# Branch: direct vs rag
add_rounded_rect(slide, Inches(8.0), Inches(1.3), Inches(1.5), Inches(0.5), ACCENT3)
add_text(slide, Inches(8.0), Inches(1.3), Inches(1.5), Inches(0.5),
         "route_question", font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)

# direct branch
add_rounded_rect(slide, Inches(10.0), Inches(1.1), Inches(1.8), Inches(0.5), ORANGE)
add_text(slide, Inches(10.0), Inches(1.1), Inches(1.8), Inches(0.5),
         "generate_direct_answer", font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)
add_arrow(slide, Inches(9.5), Inches(1.55), Inches(10.0), Inches(1.35), ORANGE, 1.5)

# rag branch - retrieve
add_rounded_rect(slide, Inches(8.0), Inches(2.5), Inches(1.5), Inches(0.5), ACCENT)
add_text(slide, Inches(8.0), Inches(2.5), Inches(1.5), Inches(0.5),
         "retrieve", font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)
add_arrow(slide, Inches(8.75), Inches(1.8), Inches(8.75), Inches(2.5), ACCENT, 1.5)

# check quality
add_rounded_rect(slide, Inches(8.0), Inches(3.3), Inches(1.5), Inches(0.5), ACCENT)
add_text(slide, Inches(8.0), Inches(3.3), Inches(1.5), Inches(0.5),
         "check_retrieval_quality", font_size=8, color=WHITE, alignment=PP_ALIGN.CENTER)
add_arrow(slide, Inches(8.75), Inches(3.0), Inches(8.75), Inches(3.3), ACCENT, 1.5)

# good -> generate
add_rounded_rect(slide, Inches(10.0), Inches(3.1), Inches(1.8), Inches(0.5), ACCENT3)
add_text(slide, Inches(10.0), Inches(3.1), Inches(1.8), Inches(0.5),
         "generate_rag_answer", font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)
add_arrow(slide, Inches(9.5), Inches(3.55), Inches(10.0), Inches(3.35), ACCENT3, 1.5)

# poor -> rewrite
add_rounded_rect(slide, Inches(5.8), Inches(3.3), Inches(1.5), Inches(0.5), PINK)
add_text(slide, Inches(5.8), Inches(3.3), Inches(1.5), Inches(0.5),
         "rewrite_question", font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)
add_arrow(slide, Inches(8.0), Inches(3.55), Inches(7.3), Inches(3.55), PINK, 1.5)

# retrieve_rewritten
add_rounded_rect(slide, Inches(3.8), Inches(3.3), Inches(1.5), Inches(0.5), PINK)
add_text(slide, Inches(3.8), Inches(3.3), Inches(1.5), Inches(0.5),
         "retrieve_rewritten", font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)
add_arrow(slide, Inches(5.8), Inches(3.55), Inches(5.3), Inches(3.55), PINK, 1.5)

# check_rewritten_quality
add_rounded_rect(slide, Inches(1.8), Inches(3.3), Inches(1.5), Inches(0.5), PINK)
add_text(slide, Inches(1.8), Inches(3.3), Inches(1.5), Inches(0.5),
         "check_rewritten_quality", font_size=8, color=WHITE, alignment=PP_ALIGN.CENTER)
add_arrow(slide, Inches(3.8), Inches(3.55), Inches(3.3), Inches(3.55), PINK, 1.5)

# rewritten_good -> generate
add_rounded_rect(slide, Inches(1.8), Inches(4.2), Inches(1.8), Inches(0.5), ACCENT3)
add_text(slide, Inches(1.8), Inches(4.2), Inches(1.8), Inches(0.5),
         "generate_rag_answer", font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)
add_arrow(slide, Inches(2.55), Inches(3.8), Inches(2.55), Inches(4.2), ACCENT3, 1.5)

# rewritten_poor -> no_context
add_rounded_rect(slide, Inches(4.0), Inches(4.2), Inches(2.0), Inches(0.5), ORANGE)
add_text(slide, Inches(4.0), Inches(4.2), Inches(2.0), Inches(0.5),
         "generate_no_context_answer", font_size=8, color=WHITE, alignment=PP_ALIGN.CENTER)
add_arrow(slide, Inches(3.3), Inches(4.45), Inches(4.0), Inches(4.45), ORANGE, 1.5)

# END nodes
add_rounded_rect(slide, Inches(10.5), Inches(4.2), Inches(1.0), Inches(0.4), RGBColor(0x37, 0x41, 0x51))
add_text(slide, Inches(10.5), Inches(4.2), Inches(1.0), Inches(0.4),
         "END", font_size=10, bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)
add_arrow(slide, Inches(10.9), Inches(3.6), Inches(10.9), Inches(4.2), GRAY, 1)
add_arrow(slide, Inches(10.0), Inches(1.35), Inches(10.5), Inches(4.2), GRAY, 1)
add_arrow(slide, Inches(6.0), Inches(4.45), Inches(10.5), Inches(4.45), GRAY, 1)

# --- Right side: session memory ---
add_rounded_rect(slide, Inches(0.3), Inches(5.2), Inches(12.7), Inches(2.0), BG_CARD)
add_rect(slide, Inches(0.3), Inches(5.2), Inches(0.08), Inches(2.0), ACCENT)
add_text(slide, Inches(0.5), Inches(5.35), Inches(12), Inches(0.35),
         "会话记忆策略", font_size=14, bold=True, color=ACCENT)

memory_details = [
    ("≤ 16 条消息", "使用完整历史"),
    ("> 16 条消息", "生成会话摘要，后续使用 摘要 + 最近 8 条原始消息"),
    ("摘要上限", "约 600 字，降低长对话 token 消耗"),
    ("上下文补全", "complete_question_with_history 节点将指代不明的问题改写为独立问题"),
]

y_off = Inches(5.8)
for title, desc in memory_details:
    add_text(slide, Inches(0.5), y_off, Inches(2.5), Inches(0.3),
             title, font_size=11, bold=True, color=WHITE)
    add_text(slide, Inches(3.0), y_off, Inches(9.5), Inches(0.3),
             desc, font_size=10, color=GRAY)
    y_off += Inches(0.38)


# ============================================================
# Slide 6: 核心流程 - 文档入库
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_page_header(slide, "05  核心流程 — 文档入库", 6, 12)

add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
         "Redis + Celery 异步入库 Pipeline", font_size=20, bold=True, color=WHITE)

# Pipeline flow
pipeline_nodes = [
    ("用户上传", "前端上传文档\n支持 txt/md/pdf/docx\nOffice/图片/HTML", ACCENT, Inches(0.3), Inches(1.6)),
    ("创建 Task", "写入 pending 文档\n创建 IngestionTask\n状态: pending", ACCENT2, Inches(3.0), Inches(1.6)),
    ("Celery 排队", "Redis Broker 分发\nworker 消费任务\n支持排队/重试/取消", ACCENT3, Inches(5.7), Inches(1.6)),
    ("inspect_document", "检测文档类型\n提取文本内容\n获取存储元数据", ORANGE, Inches(8.4), Inches(1.6)),
    ("chunk_embed_index", "分块(800字符)\nEmbedding 向量化\n写入 Milvus", PINK, Inches(0.3), Inches(3.8)),
    ("update_document_record", "更新文档记录\n同步 chunk 元数据\n记录节点日志", ACCENT, Inches(3.0), Inches(3.8)),
    ("完成", "状态: success\n管理后台可查看\n支持重试/取消", ACCENT3, Inches(5.7), Inches(3.8)),
]

for i, (name, desc, color, x, y) in enumerate(pipeline_nodes):
    w = Inches(2.4)
    h = Inches(1.5)
    add_rounded_rect(slide, x, y, w, h, BG_CARD)
    add_rect(slide, x, y, w, Pt(3), color)
    add_text(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.3),
             name, font_size=11, bold=True, color=color)
    add_text(slide, x + Inches(0.1), y + Inches(0.35), w - Inches(0.2), Inches(1.1),
             desc, font_size=9, color=GRAY)

# Arrows between pipeline nodes
arrow_pairs = [
    (0, 1), (1, 2), (2, 3), (3, 4), (4, 5), (5, 6)
]
for a, b in arrow_pairs:
    x1 = pipeline_nodes[a][3] + Inches(2.4)
    y1 = pipeline_nodes[a][4] + Inches(0.75)
    x2 = pipeline_nodes[b][3]
    y2 = pipeline_nodes[b][4] + Inches(0.75)
    add_arrow(slide, x1, y1, x2, y2, GRAY, 1.5)

# --- Bottom: retry & cancel ---
add_rounded_rect(slide, Inches(0.3), Inches(5.8), Inches(12.7), Inches(1.4), BG_CARD)
add_rect(slide, Inches(0.3), Inches(5.8), Inches(0.08), Inches(1.4), ORANGE)
add_text(slide, Inches(0.5), Inches(5.95), Inches(12), Inches(0.35),
         "任务管理", font_size=14, bold=True, color=ORANGE)

task_mgmt = [
    ("重试机制", "失败后自动重试，默认最多 3 次，间隔 30s；认证错误等不可重试错误直接标记 failed"),
    ("取消操作", "支持取消 pending / queued / running / retrying 状态的任务"),
    ("节点日志", "每个 Pipeline 节点记录状态、耗时、chunks 数量、跳过数量和错误信息"),
    ("解析器选择", "支持 Docling / Azure Document Intelligence / basic 三种文档解析路径"),
]

y_off = Inches(6.4)
for title, desc in task_mgmt:
    add_text(slide, Inches(0.5), y_off, Inches(2.5), Inches(0.3),
             title, font_size=11, bold=True, color=WHITE)
    add_text(slide, Inches(3.0), y_off, Inches(9.5), Inches(0.3),
             desc, font_size=9, color=GRAY)
    y_off += Inches(0.32)


# ============================================================
# Slide 7: 核心代码 - LangGraph Agent
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_page_header(slide, "06  核心代码 — LangGraph Agent 编排", 7, 12)

add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
         "graph_service.py — 构建 Agent 工作流图", font_size=18, bold=True, color=WHITE)

# Left: code
code_lines = [
    "def build_chat_graph():",
    "    graph = StateGraph(ChatState)",
    "    # 添加节点",
    "    graph.add_node('complete_question_with_history', ...)",
    "    graph.add_node('analyze_question', ...)",
    "    graph.add_node('plan_agent_task', ...)",
    "    graph.add_node('retrieve', ...)",
    "    graph.add_node('check_retrieval_quality', ...)",
    "    graph.add_node('rewrite_question', ...)",
    "    graph.add_node('retrieve_rewritten', ...)",
    "    graph.add_node('generate_rag_answer', ...)",
    "    graph.add_node('generate_direct_answer', ...)",
    "    # 添加边",
    "    graph.add_edge(START, 'complete_question_with_history')",
    "    graph.add_conditional_edges(",
    "        'plan_agent_task', route_question,",
    "        {'rag': 'retrieve', 'direct': 'generate_direct_answer'}",
    "    )",
    "    graph.add_conditional_edges(",
    "        'check_retrieval_quality', route_retrieval_quality,",
    "        {'good': 'generate_rag_answer', 'poor': 'rewrite_question'}",
    "    )",
    "    return graph.compile()",
]

add_code_block(slide, Inches(0.3), Inches(1.4), Inches(6.5), Inches(5.8), code_lines, font_size=9.5)

# Right: ChatState
add_rounded_rect(slide, Inches(7.1), Inches(1.4), Inches(5.9), Inches(5.8), BG_CARD)
add_rect(slide, Inches(7.1), Inches(1.4), Inches(0.08), Inches(5.8), ACCENT2)
add_text(slide, Inches(7.3), Inches(1.55), Inches(5.5), Inches(0.35),
         "ChatState 状态定义", font_size=14, bold=True, color=ACCENT2)

state_fields = [
    ("question", "用户原始问题"),
    ("chat_history", "会话历史文本"),
    ("standalone_question", "上下文补全后的独立问题"),
    ("rewritten_question", "检索改写后的查询"),
    ("route", "路由结果: rag / direct"),
    ("task_intent", "任务意图标签"),
    ("sources", "检索结果列表"),
    ("retrieval_quality", "检索质量: good / poor"),
    ("agent_plan", "Agent 执行计划"),
    ("tool_calls", "工具调用记录"),
    ("answer", "最终回答"),
    ("steps", "执行步骤日志"),
]

y_off = Inches(2.0)
for field, desc in state_fields:
    add_text(slide, Inches(7.3), y_off, Inches(1.5), Inches(0.28),
             field, font_size=9, bold=True, color=ACCENT)
    add_text(slide, Inches(8.8), y_off, Inches(4.0), Inches(0.28),
             desc, font_size=9, color=GRAY)
    y_off += Inches(0.4)


# ============================================================
# Slide 8: 核心代码 - 检索与入库
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_page_header(slide, "07  核心代码 — 检索 & 入库", 8, 12)

add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
         "rag_service.py & ingestion_pipeline.py", font_size=18, bold=True, color=WHITE)

# Left: retrieval code
add_rounded_rect(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.9), BG_CARD)
add_rect(slide, Inches(0.3), Inches(1.3), Inches(0.08), Inches(5.9), ACCENT)
add_text(slide, Inches(0.5), Inches(1.45), Inches(5.8), Inches(0.3),
         "混合检索 + 质量评分", font_size=13, bold=True, color=ACCENT)

retrieval_code = [
    "def retrieve_sources_multi(collection_names, question, top_k=4):",
    "    # 1. 向量检索 (COSINE)",
    "    vector_results = milvus.search(data=[query_vector], ...)",
    "    # 2. 关键词检索 (jieba + ILIKE)",
    "    keyword_results = db.query(...).filter(or_(*filters))",
    "    # 3. 混合评分融合",
    "    for source in merged_sources:",
    "        key = (source.source, source.content)",
    "        existing = deduplicated_by_chunk.get(key)",
    "        if existing:",
    "            merge_source_scores(existing, source)  # 去重+取最高分",
    "    # 4. 最终排序",
    "    return sorted_sources[:top_k]",
    "",
    "def calculate_retrieval_score(source):",
    "    if vector_score > 0 and keyword_score > 0:",
    "        return vector_score * 0.7 + keyword_score * 0.3 + 0.05",
    "    if vector_score > 0:",
    "        return vector_score * 0.7",
    "    return keyword_score * 0.3",
]

add_code_block(slide, Inches(0.4), Inches(1.85), Inches(6.0), Inches(5.3), retrieval_code, font_size=9)

# Right: ingestion pipeline code
add_rounded_rect(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.9), BG_CARD)
add_rect(slide, Inches(6.8), Inches(1.3), Inches(0.08), Inches(5.9), ACCENT3)
add_text(slide, Inches(7.0), Inches(1.45), Inches(5.8), Inches(0.3),
         "入库 Pipeline 节点", font_size=13, bold=True, color=ACCENT3)

pipeline_code = [
    "def run_ingestion_pipeline(db, task):",
    "    # 节点1: 检查文档",
    "    with task_node(db, task, 'inspect_document'):",
    "        metadata = get_stored_file_metadata(source)",
    "        details.update({...})",
    "",
    "    # 节点2: 分块 + 向量化 + 入库",
    "    with task_node(db, task, 'chunk_embed_index'):",
    "        chunks, skipped = ingest_document(",
    "            collection_name, source,",
    "            embedding_model=kb.embedding_model,",
    "            db=db, knowledge_base_id=kb.id,",
    "        )",
    "        details.update({'chunks': chunks, 'skipped': skipped})",
    "",
    "    # 节点3: 更新文档记录",
    "    with task_node(db, task, 'update_document_record'):",
    "        upsert_document_record(db, ...)",
    "",
    "    return update_ingestion_task(db, task, status='success')",
]

add_code_block(slide, Inches(6.9), Inches(1.85), Inches(6.0), Inches(5.3), pipeline_code, font_size=9)


# ============================================================
# Slide 9: 核心代码 - 可观测性
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_page_header(slide, "08  核心代码 — 可观测性 & 配置", 9, 12)

add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
         "Langfuse 链路追踪 & Pydantic Settings 配置", font_size=18, bold=True, color=WHITE)

# Left: observability
add_rounded_rect(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.9), BG_CARD)
add_rect(slide, Inches(0.3), Inches(1.3), Inches(0.08), Inches(5.9), PINK)
add_text(slide, Inches(0.5), Inches(1.45), Inches(5.8), Inches(0.3),
         "observability.py — 节点级追踪", font_size=13, bold=True, color=PINK)

obs_code = [
    "def traced_node(name, node_func):",
    "    def wrapped(state):",
    "        with start_node_span(",
    "            name,",
    "            input_data=state_snapshot_for_span(state)",
    "        ) as span:",
    "            result = node_func(state)",
    "            update_node_span(span,",
    "                output=serialize_value_for_span(result),",
    "                metadata={'node': name})",
    "            return result",
    "    return wrapped",
    "",
    "# 每个 Agent 节点自动记录:",
    "# - 输入: question, route, sources, steps",
    "# - 输出: 节点返回结果",
    "# - 模型调用: input/output/usage/token",
    "# - Langfuse 可视化全链路",
]

add_code_block(slide, Inches(0.4), Inches(1.85), Inches(6.0), Inches(5.3), obs_code, font_size=9)

# Right: config
add_rounded_rect(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.9), BG_CARD)
add_rect(slide, Inches(6.8), Inches(1.3), Inches(0.08), Inches(5.9), ACCENT2)
add_text(slide, Inches(7.0), Inches(1.45), Inches(5.8), Inches(0.3),
         "config.py — 环境变量配置", font_size=13, bold=True, color=ACCENT2)

config_code = [
    "class Settings(BaseSettings):",
    "    # 模型",
    "    openai_api_key: str",
    "    chat_model: str = 'qwen3.7-plus'",
    "    embedding_model: str = 'qwen3.7-text-embedding'",
    "    # 存储",
    "    milvus_uri: str",
    "    database_url: str",
    "    s3_endpoint_url: str = 'http://127.0.0.1:9000'",
    "    # 认证",
    "    auth_secret_key: str",
    "    seed_admin_password: str = 'admin123456'",
    "    # 会话记忆",
    "    chat_memory_summary_max_chars: int = 600",
    "    chat_memory_summary_start_messages: int = 16",
    "    # Rerank",
    "    rerank_enabled: bool = False",
    "    rerank_model: str = 'qwen3-rerank'",
    "    # 文档解析",
    "    document_parser: str = 'docling'",
    "    docling_do_ocr: bool = True",
    "",
    "    model_config = SettingsConfigDict(env_file='.env')",
]

add_code_block(slide, Inches(6.9), Inches(1.85), Inches(6.0), Inches(5.3), config_code, font_size=9)


# ============================================================
# Slide 10: 部署说明
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_page_header(slide, "09  部署说明", 10, 12)

add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
         "Docker Compose 一键部署", font_size=20, bold=True, color=WHITE)

# Left: deployment steps
add_rounded_rect(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.9), BG_CARD)
add_rect(slide, Inches(0.3), Inches(1.3), Inches(0.08), Inches(5.9), ACCENT)
add_text(slide, Inches(0.5), Inches(1.45), Inches(5.8), Inches(0.3),
         "部署步骤", font_size=14, bold=True, color=ACCENT)

steps = [
    ("1. 准备环境变量", "cp .env.example .env\n至少填写 OPENAI_API_KEY"),
    ("2. 一键启动", "docker compose up -d --build"),
    ("3. 查看状态", "docker compose ps\ndocker compose logs -f backend"),
    ("4. 访问服务", "前端: http://localhost:3000\n后端: http://127.0.0.1:8000/docs\nMinIO: http://127.0.0.1:9001"),
    ("5. 健康检查", "GET /health\n返回: { status: 'ok', services: {...} }"),
]

y_off = Inches(1.9)
for title, desc in steps:
    add_text(slide, Inches(0.5), y_off, Inches(5.8), Inches(0.3),
             title, font_size=11, bold=True, color=WHITE)
    for i, line in enumerate(desc.split("\n")):
        add_text(slide, Inches(0.7), y_off + Inches(0.3) + Inches(0.25) * i,
                 Inches(5.5), Inches(0.25),
                 line, font_size=9, color=GRAY)
    y_off += Inches(0.95)

# Right: services table
add_rounded_rect(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(3.5), BG_CARD)
add_rect(slide, Inches(6.8), Inches(1.3), Inches(0.08), Inches(3.5), ACCENT3)
add_text(slide, Inches(7.0), Inches(1.45), Inches(5.8), Inches(0.3),
         "服务端口映射", font_size=14, bold=True, color=ACCENT3)

services_table = [
    ("frontend", "Next.js 前端", "3000"),
    ("backend", "FastAPI 后端", "8000"),
    ("postgres", "PostgreSQL", "5433 → 5432"),
    ("milvus", "Milvus 向量库", "19530"),
    ("minio", "MinIO 对象存储", "9000 / 9001"),
    ("redis", "Redis", "6379"),
    ("docling", "Docling 解析", "5001"),
]

y_off = Inches(1.9)
for name, desc, port in services_table:
    add_text(slide, Inches(7.0), y_off, Inches(1.5), Inches(0.28),
             name, font_size=10, bold=True, color=ACCENT)
    add_text(slide, Inches(8.5), y_off, Inches(2.5), Inches(0.28),
             desc, font_size=10, color=WHITE)
    add_text(slide, Inches(11.0), y_off, Inches(1.8), Inches(0.28),
             port, font_size=10, color=GRAY)
    y_off += Inches(0.38)

# Bottom: data volumes & backup
add_rounded_rect(slide, Inches(6.8), Inches(5.0), Inches(6.2), Inches(2.2), BG_CARD)
add_rect(slide, Inches(6.8), Inches(5.0), Inches(0.08), Inches(2.2), ORANGE)
add_text(slide, Inches(7.0), Inches(5.15), Inches(5.8), Inches(0.3),
         "数据持久化 & 备份", font_size=14, bold=True, color=ORANGE)

backup_info = [
    ("数据卷", "postgres_data / milvus_data / minio_data / etcd_data / redis_data"),
    ("备份建议", "定期导出 PostgreSQL + MinIO bucket + Milvus 数据卷快照"),
    ("注意事项", "不要执行 docker compose down -v（会删除所有数据卷）"),
]

y_off = Inches(5.55)
for title, desc in backup_info:
    add_text(slide, Inches(7.0), y_off, Inches(2.0), Inches(0.3),
             title, font_size=10, bold=True, color=WHITE)
    add_text(slide, Inches(9.0), y_off, Inches(3.8), Inches(0.3),
             desc, font_size=9, color=GRAY)
    y_off += Inches(0.45)


# ============================================================
# Slide 11: 项目结构
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_page_header(slide, "10  项目结构", 11, 12)

add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
         "代码组织", font_size=20, bold=True, color=WHITE)

# Left: backend structure
add_rounded_rect(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.9), BG_CARD)
add_rect(slide, Inches(0.3), Inches(1.3), Inches(0.08), Inches(5.9), ACCENT)
add_text(slide, Inches(0.5), Inches(1.45), Inches(5.8), Inches(0.3),
         "rag-backend/app/", font_size=14, bold=True, color=ACCENT)

backend_structure = [
    ("api/routes/", "FastAPI 路由层"),
    ("  auth.py", "登录 / 获取当前用户"),
    ("  chat.py", "聊天问答接口"),
    ("  knowledge_bases.py", "知识库 CRUD"),
    ("  documents.py", "文档管理"),
    ("  sessions.py", "会话管理"),
    ("  ingestion/", "入库任务接口"),
    ("core/", "核心初始化"),
    ("  config.py", "Pydantic Settings"),
    ("  db.py", "SQLAlchemy 引擎"),
    ("  celery_app.py", "Celery 实例"),
    ("models/", "SQLAlchemy 数据模型"),
    ("  user.py / chat.py / document.py", "用户/消息/文档"),
    ("  knowledge_base.py", "知识库模型"),
    ("  ingestion_task.py", "入库任务模型"),
    ("schemas/", "请求/响应 Pydantic 模型"),
    ("services/", "业务逻辑层"),
    ("  graph_service.py", "LangGraph Agent 编排"),
    ("  rag_service.py", "检索 + 入库核心"),
    ("  ingestion_pipeline.py", "入库 Pipeline"),
    ("  storage_service.py", "MinIO 存储"),
    ("  auth_service.py", "JWT 认证"),
    ("  observability.py", "Langfuse 追踪"),
    ("tasks/", "Celery 异步任务"),
    ("  ingestion.py", "入库任务执行器"),
]

y_off = Inches(1.85)
for path, desc in backend_structure:
    is_dir = path.endswith("/")
    add_text(slide, Inches(0.5), y_off, Inches(2.5), Inches(0.28),
             path, font_size=9, bold=is_dir, color=ACCENT if is_dir else LIGHT_GRAY)
    add_text(slide, Inches(3.0), y_off, Inches(3.2), Inches(0.28),
             desc, font_size=9, color=GRAY)
    y_off += Inches(0.32)

# Right: frontend structure
add_rounded_rect(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.9), BG_CARD)
add_rect(slide, Inches(6.8), Inches(1.3), Inches(0.08), Inches(5.9), ACCENT2)
add_text(slide, Inches(7.0), Inches(1.45), Inches(5.8), Inches(0.3),
         "rag-frontend/", font_size=14, bold=True, color=ACCENT2)

frontend_structure = [
    ("app/", "Next.js App Router 页面"),
    ("  page.tsx", "聊天主页面"),
    ("  page.module.css", "页面样式"),
    ("  layout.tsx", "根布局"),
    ("  globals.css", "全局样式"),
    ("  login/", "登录页面"),
    ("  admin/", "管理后台页面"),
    ("components/", "公共组件"),
    ("  Modal.tsx", "模态框"),
    ("  Select.tsx", "选择器"),
    ("  Pagination/", "分页组件"),
    ("  MenuIcon.tsx", "菜单图标"),
    ("lib/", "前端工具库"),
    ("  api.ts", "API 请求封装"),
    ("  auth.ts", "认证工具"),
    ("", ""),
    ("关键技术点", ""),
    ("  • Next.js 16 App Router", ""),
    ("  • React 19 + TypeScript", ""),
    ("  • 流式 SSE 渲染", ""),
    ("  • Radix UI 组件", ""),
    ("  • react-markdown 渲染", ""),
]

y_off = Inches(1.85)
for path, desc in frontend_structure:
    if path and path[0].isupper() and "/" not in path:
        # section header
        add_text(slide, Inches(7.0), y_off, Inches(5.5), Inches(0.3),
                 desc, font_size=11, bold=True, color=ACCENT2)
    elif path:
        add_text(slide, Inches(7.0), y_off, Inches(2.5), Inches(0.28),
                 path, font_size=9, bold="/" in path and "." not in path,
                 color=ACCENT2 if "/" in path and "." not in path else LIGHT_GRAY)
        add_text(slide, Inches(9.5), y_off, Inches(3.3), Inches(0.28),
                 desc, font_size=9, color=GRAY)
    y_off += Inches(0.32)


# ============================================================
# Slide 12: 总结与展望
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_page_header(slide, "11  总结与展望", 12, 12)

add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
         "项目总结与未来规划", font_size=20, bold=True, color=WHITE)

# Top: achievements
add_rounded_rect(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.8), BG_CARD)
add_rect(slide, Inches(0.3), Inches(1.3), Inches(0.08), Inches(2.8), ACCENT3)
add_text(slide, Inches(0.5), Inches(1.45), Inches(5.8), Inches(0.3),
         "已实现能力", font_size=14, bold=True, color=ACCENT3)

achievements = [
    "✅ Agentic RAG 多节点工作流（LangGraph）",
    "✅ 多知识库支持 + 混合检索（向量 + 关键词）",
    "✅ 会话记忆与上下文补全",
    "✅ 异步入库 Pipeline（Celery + Redis）",
    "✅ 全链路可观测性（Langfuse）",
    "✅ 多用户认证 + 管理后台",
    "✅ Docker Compose 一键部署",
    "✅ 流式输出 + 检索质量检查 + 自动改写",
]

y_off = Inches(1.85)
for ach in achievements:
    add_text(slide, Inches(0.5), y_off, Inches(5.8), Inches(0.28),
             ach, font_size=10, color=LIGHT_GRAY)
    y_off += Inches(0.3)

# Top right: future plans
add_rounded_rect(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.8), BG_CARD)
add_rect(slide, Inches(6.8), Inches(1.3), Inches(0.08), Inches(2.8), ACCENT2)
add_text(slide, Inches(7.0), Inches(1.45), Inches(5.8), Inches(0.3),
         "未来规划", font_size=14, bold=True, color=ACCENT2)

future_plans = [
    "🔜 多模态支持（图片理解、表格识别增强）",
    "🔜 Agent 工具调用（外部 API / 数据库查询）",
    "🔜 更精细的权限管理（RBAC）",
    "🔜 知识库协作与版本管理",
    "🔜 检索评估数据集与自动化评测",
    "🔜 支持更多 LLM 提供商（本地模型）",
    "🔜 前端 UI/UX 优化与移动端适配",
    "🔜 Kubernetes 部署方案",
]

y_off = Inches(1.85)
for plan in future_plans:
    add_text(slide, Inches(7.0), y_off, Inches(5.8), Inches(0.28),
             plan, font_size=10, color=LIGHT_GRAY)
    y_off += Inches(0.3)

# Bottom: key metrics / highlights
add_rounded_rect(slide, Inches(0.3), Inches(4.4), Inches(12.7), Inches(2.8), BG_CARD)
add_rect(slide, Inches(0.3), Inches(4.4), Inches(0.08), Inches(2.8), ACCENT)
add_text(slide, Inches(0.5), Inches(4.55), Inches(12), Inches(0.3),
         "技术亮点", font_size=14, bold=True, color=ACCENT)

highlights = [
    ("Agent 编排", "LangGraph 状态机驱动\n12 个节点\n条件边路由", ACCENT),
    ("检索策略", "向量 + 关键词混合\n自动改写重试\n质量检查门控", ACCENT2),
    ("异步架构", "Celery + Redis\nPipeline 节点日志\n重试/取消机制", ACCENT3),
    ("可观测性", "Langfuse 全链路\n节点级 span\nToken 用量追踪", ORANGE),
    ("部署便捷", "Docker Compose\n一键启动\n健康检查", PINK),
]

x_off = Inches(0.5)
for title, desc, color in highlights:
    w = Inches(2.4)
    add_rect(slide, x_off, Inches(5.0), w, Pt(2), color)
    add_text(slide, x_off, Inches(5.05), w, Inches(0.3),
             title, font_size=11, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text(slide, x_off, Inches(5.35), w, Inches(1.3),
             desc, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)
    x_off += Inches(2.5)

# Footer
add_text(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.3),
         "myAgent — 本地优先的 RAG Agent 系统  |  FastAPI · LangGraph · Milvus · Next.js",
         font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)


# ============ Save ============
output_path = r"D:\code\git_localRepository\myAgent\docs\ppt\myAgent_技术文档.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
